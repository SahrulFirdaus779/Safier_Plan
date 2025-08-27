import streamlit as st
import pandas as pd
import time
from datetime import datetime, date, timedelta
import uuid
from itertools import groupby
import sqlite3
from streamlit_calendar import calendar
import plotly.express as px
import json
import os
from pptx import Presentation # --- BARU ---
from pptx.util import Inches # --- BARU ---
from io import BytesIO # --- BARU ---

# --- KONFIGURASI APLIKASI ---
st.set_page_config(
    page_title="Safier Plan",
    page_icon=" briefcase ",
    layout="wide"
)

# --- CSS Kustom ---
st.markdown("""
    <style>
        /* (CSS tidak berubah dari versi sebelumnya) */
        .main, [data-testid="stAppViewContainer"] { background-color: #FFF7E8; font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; }
        .stButton > button { background-color: #007BFF !important; color: white !important; font-weight: 600 !important; border-radius: 8px !important; border: none !important; padding: 0.5rem 1rem !important; margin-top: 0.5rem !important; transition: background-color 0.3s, transform 0.2s; }
        .stButton > button:hover { background-color: #0056b3 !important; transform: translateY(-2px); }
        [data-testid="stHorizontalBlock"] .stButton > button { padding: 0.25rem 0.75rem !important; font-size: 0.85rem !important; }
        button[kind="primary"] { background-color: #DC3545 !important; }
        button[kind="primary"]:hover { background-color: #c82333 !important; }
        [data-testid="stForm"], div[data-testid="stExpander"], .st-container[border="true"] { background-color: #FFFFFF !important; border: 1px solid #DEE2E6 !important; border-radius: 12px !important; padding: 1.5rem !important; box-shadow: 0 4px 12px rgba(0,0,0,0.05); }
        [data-testid="stDateInput"] input, [data-testid="stTextInput"] input, [data-testid="stNumberInput"] input { border-radius: 8px; border: 1px solid #ced4da; padding: 0.75rem; transition: border-color 0.2s, box-shadow 0.2s; }
        [data-testid="stDateInput"] input:focus, [data-testid="stTextInput"] input:focus, [data-testid="stNumberInput"] input:focus { border-color: #007BFF; box-shadow: 0 0 0 0.2rem rgba(0,123,255,.25); }
        div[data-baseweb="tab-list"] { gap: 8px; }
        button[data-baseweb="tab"] { background-color: transparent; font-size: 1.1rem; font-weight: 600; border-radius: 8px 8px 0 0 !important; border-bottom: 2px solid transparent !important; padding: 0.5rem 1rem; }
        button[data-baseweb="tab"][aria-selected="true"] { background-color: #007BFF !important; color: white !important; border-bottom: 2px solid #0056b3 !important; }
        .table-header { background-color: #007BFF; color: white; padding: 0.75rem; border-radius: 8px; font-weight: 600; text-align: center; }
        .status-badge { padding: 0.25rem 0.6rem; border-radius: 12px; font-size: 0.8rem; font-weight: 600; color: white; text-align: center; }
        .status-completed { background-color: #28a745; }
        .status-inprogress { background-color: #fd7e14; }
        .status-scheduled { background-color: #17a2b8; }
        .status-notstarted { background-color: #6c757d; }
    </style>
    """, unsafe_allow_html=True)

# --- FUNGSI-FUNGSI DATABASE (Tidak Berubah) ---
DB_FILE = "tasks.db"
DATE_COLUMNS = ["Tanggal Mulai", "Tanggal Selesai Target", "Tanggal Jadwal", "Tanggal Selesai"]
def db_connect(): return sqlite3.connect(DB_FILE)
def setup_database():
    conn = db_connect()
    cursor = conn.cursor()
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS tasks (
            id TEXT PRIMARY KEY, Tugas TEXT NOT NULL, Deskripsi TEXT,
            "Durasi (jam)" REAL, "Tanggal Mulai" TEXT, "Tanggal Selesai Target" TEXT,
            Selesai BOOLEAN, Prioritas TEXT, Delegasi TEXT,
            "Tanggal Jadwal" TEXT, "Tanggal Selesai" TEXT, SubTugas TEXT
        )
    """)
    conn.commit()
    conn.close()
def load_tasks_from_db():
    if not os.path.exists(DB_FILE): return []
    conn = db_connect()
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM tasks")
    tasks_rows = cursor.fetchall()
    conn.close()
    tasks = []
    for row in tasks_rows:
        task = dict(row)
        for date_col in DATE_COLUMNS:
            if task.get(date_col):
                try: task[date_col] = datetime.strptime(task[date_col], '%Y-%m-%d').date()
                except (ValueError, TypeError): task[date_col] = None
        task['SubTugas'] = json.loads(task.get('SubTugas', '[]'))
        tasks.append(task)
    return tasks
def save_task_to_db(task_dict):
    conn = db_connect()
    cursor = conn.cursor()
    task_to_save = task_dict.copy()
    for date_col in DATE_COLUMNS:
        if isinstance(task_to_save.get(date_col), date):
            task_to_save[date_col] = task_to_save[date_col].isoformat()
    if 'SubTugas' in task_to_save:
        task_to_save['SubTugas'] = json.dumps(task_to_save['SubTugas'])
    cursor.execute("PRAGMA table_info(tasks)")
    table_columns = [info[1] for info in cursor.fetchall()]
    task_to_save = {k: v for k, v in task_to_save.items() if k in table_columns}
    columns = ', '.join([f'"{key}"' for key in task_to_save.keys()])
    placeholders = ', '.join(['?'] * len(task_to_save))
    sql = f"INSERT OR REPLACE INTO tasks ({columns}) VALUES ({placeholders})"
    cursor.execute(sql, list(task_to_save.values()))
    conn.commit()
    conn.close()
def delete_task_from_db(task_id):
    conn = db_connect()
    cursor = conn.cursor()
    cursor.execute("DELETE FROM tasks WHERE id = ?", (task_id,))
    conn.commit()
    conn.close()

# --- INISIALISASI & LOGIKA UTAMA (Tidak Berubah) ---
setup_database()
if 'tasks' not in st.session_state: st.session_state.tasks = load_tasks_from_db()
if 'editing_task_id' not in st.session_state: st.session_state.editing_task_id = None
if 'pomodoro_running' not in st.session_state: st.session_state.pomodoro_running = False
if 'active_pomodoro_task' not in st.session_state: st.session_state.active_pomodoro_task = None
if 'pomodoro_start_time' not in st.session_state: st.session_state.pomodoro_start_time = 0
for task in st.session_state.tasks:
    if task.get('Prioritas') == 'Jadwalkan' and task.get('Tanggal Jadwal') and isinstance(task.get('Tanggal Jadwal'), date) and task.get('Tanggal Jadwal') == date.today():
        task['Prioritas'] = 'Lakukan Sekarang'
        save_task_to_db(task)
        st.toast(f"Tugas '{task.get('Tugas')}' kini menjadi 'Lakukan Sekarang'!", icon="✨")
def get_task_by_id(task_id): return next((task for task in st.session_state.tasks if task.get('id') == task_id), None)
def update_task_priority_by_id(task_id, new_priority):
    task = get_task_by_id(task_id)
    if task: 
        task['Prioritas'] = new_priority
        save_task_to_db(task)
def get_task_status(task):
    if task.get('Selesai'): return "Completed", "status-completed", "#28a745"
    if task.get('Prioritas') == 'Lakukan Sekarang': return "In Progress", "status-inprogress", "#fd7e14"
    if task.get('Prioritas') == 'Jadwalkan': return "Scheduled", "status-scheduled", "#17a2b8"
    return "Not Started", "status-notstarted", "#6c757d"

# --- (BARU) FUNGSI UNTUK MEMBUAT PRESENTASI ---
def buat_presentasi_laporan(tasks):
    prs = Presentation()
    
    # Slide 1: Judul
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Laporan Produktivitas - Safier Plan"
    subtitle.text = f"Dibuat pada: {date.today().strftime('%d %B %Y')}"

    # Slide 2: Ringkasan Metrik
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Ringkasan Produktivitas"
    
    tasks_aktif = [t for t in tasks if not t.get('Selesai')]
    tasks_selesai_hari_ini = [t for t in tasks if t.get('Selesai') and t.get('Tanggal Selesai') == date.today()]
    tasks_mendesak = [t for t in tasks_aktif if t.get('Prioritas') == 'Lakukan Sekarang']
    
    body_shape = content.text_frame
    body_shape.clear()
    p = body_shape.paragraphs[0]
    p.text = f"Tugas Selesai Hari Ini: {len(tasks_selesai_hari_ini)}"
    p = body_shape.add_paragraph()
    p.text = f"Total Tugas Aktif: {len(tasks_aktif)}"
    p = body_shape.add_paragraph()
    p.text = f"Tugas Mendesak: {len(tasks_mendesak)}"

    # Slide 3: Daftar Tugas Aktif
    if tasks_aktif:
        slide_layout = prs.slide_layouts[5] # Judul saja
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Daftar Tugas Aktif"
        
        rows = len(tasks_aktif) + 1
        cols = 3
        left = Inches(1.0)
        top = Inches(1.5)
        width = Inches(8.0)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(4.0)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(2.0)

        # Header
        table.cell(0, 0).text = 'Nama Tugas'
        table.cell(0, 1).text = 'Prioritas'
        table.cell(0, 2).text = 'Tanggal Selesai Target'
        
        # Isi tabel
        for i, task in enumerate(tasks_aktif):
            table.cell(i + 1, 0).text = task.get('Tugas', '')
            table.cell(i + 1, 1).text = task.get('Prioritas', '')
            table.cell(i + 1, 2).text = task.get('Tanggal Selesai Target').strftime('%d %b %Y') if task.get('Tanggal Selesai Target') else 'N/A'

    # Simpan presentasi ke dalam memori
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# --- TAMPILAN UTAMA APLIKASI ---
st.title("Safier Plan")
st.markdown("Sistem terintegrasi untuk mengelola waktu dan tugas Anda secara efektif.")
st.header("Ringkasan Produktivitas", divider='rainbow')
tasks_aktif = [t for t in st.session_state.tasks if not t.get('Selesai')]
tasks_selesai_hari_ini = [t for t in st.session_state.tasks if t.get('Selesai') and t.get('Tanggal Selesai') == date.today()]
tasks_mendesak = [t for t in tasks_aktif if t.get('Prioritas') == 'Lakukan Sekarang']
col1, col2, col3 = st.columns(3)
col1.metric(label="Tugas Selesai Hari Ini", value=len(tasks_selesai_hari_ini))
col2.metric(label="Total Tugas Aktif", value=len(tasks_aktif))
col3.metric(label="Tugas Mendesak", value=len(tasks_mendesak))
st.markdown("---")

# --- NAVIGASI TAB (DITAMBAH LAPORAN & EKSPOR) ---
tab_names = ["Dashboard & Input", "Prioritaskan", "Jadwal Aktivitas", "Kalender", "Sesi Fokus", "Delegasikan", "Laporan & Ekspor"]
tabs = st.tabs(tab_names)


tab1, tab2, tab3, tab4, tab5, tab6, tab7 = st.tabs(["Dashboard & Input", "Prioritaskan", "Jadwal Aktivitas", "Kalender", "Sesi Fokus", "Delegasikan", "Laporan"])

with tab1:
    st.header("Input & Kelola Semua Tugas")
    if st.session_state.editing_task_id:
        task_to_edit = get_task_by_id(st.session_state.editing_task_id)
        if task_to_edit:
            with st.form("edit_task_form"):
                st.subheader(f"Edit Tugas: {task_to_edit.get('Tugas', '')}")
                new_task_name = st.text_input("Nama Tugas", value=task_to_edit.get('Tugas', ''))
                new_desc = st.text_area("Deskripsi", value=task_to_edit.get('Deskripsi', ''))
                c1, c2 = st.columns(2)
                new_start_date = c1.date_input("Tanggal Mulai", value=task_to_edit.get('Tanggal Mulai', date.today()))
                new_end_date = c2.date_input("Tanggal Selesai", value=task_to_edit.get('Tanggal Selesai Target', date.today()))
                c1_btn, c2_btn = st.columns(2)
                if c1_btn.form_submit_button("Simpan Perubahan"):
                    task_to_edit.update({'Tugas': new_task_name, 'Deskripsi': new_desc, 'Tanggal Mulai': new_start_date, 'Tanggal Selesai Target': new_end_date})
                    save_task_to_db(task_to_edit); st.session_state.editing_task_id = None
                    st.success("Tugas berhasil diperbarui!"); st.rerun()
                if c2_btn.form_submit_button("Batal", type="secondary"): st.session_state.editing_task_id = None; st.rerun()
    with st.expander("Tambahkan Tugas Baru"):
        with st.form("new_task_form"):
            st.subheader("Detail Tugas Utama")
            task_input = st.text_input("Nama Tugas:", placeholder="Contoh: Menyelesaikan Kelas Belajar Dasar Git")
            c1, c2 = st.columns(2)
            start_date = c1.date_input("Tanggal Mulai", value=date.today())
            end_date = c2.date_input("Tanggal Selesai", value=date.today())
            st.subheader("Definisikan Tugas (opsional, dipisah koma)")
            smart_s = st.text_input("S (Specific):", placeholder="Saya akan menyelesaikan kelas Belajar Dasar Git dengan GitHub sampai dengan selesai")
            smart_m = st.text_input("M (Measurable):", placeholder="Lulus kelas Belajar Dasar Git dengan GitHub sesuai jam belajar yang ditetapkan")
            smart_a = st.text_input("A (Achievable):", placeholder="Saya memiliki kemampuan yang cukup untuk menyelesaikannya")
            smart_r = st.text_input("R (Relevant):", placeholder="Kelas ini sesuai dengan kebutuhan saat ini baik akademik maupun profesional")
            smart_t = st.text_input("T (Time-bound):", placeholder="Saya akan lulus kelas pada tanggal 30 Agustus 2025")
            if st.form_submit_button("Simpan Tugas"):
                if task_input:
                    smart_parts = [s for s in [smart_s, smart_m, smart_a, smart_r, smart_t] if s]
                    full_description = ", ".join(smart_parts) if smart_parts else "Tidak ada deskripsi."
                    new_task = {"id": str(uuid.uuid4()), "Tugas": task_input, "Deskripsi": full_description, "Durasi (jam)": 0.0, "Tanggal Mulai": start_date, "Tanggal Selesai Target": end_date, "Selesai": False, "Prioritas": "Belum Diprioritaskan", "Delegasi": "", "Tanggal Jadwal": None, "Tanggal Selesai": None}
                    st.session_state.tasks.append(new_task); save_task_to_db(new_task)
                    st.success(f"Tugas '{task_input}' berhasil ditambahkan!"); st.rerun()
                else: st.warning("Nama tugas tidak boleh kosong.")
    st.subheader("Tabel Tugas Lengkap")
    search_query = st.text_input("Cari tugas berdasarkan nama:", placeholder="Ketik nama tugas...")
    priority_options = ["Belum Diprioritaskan", "Lakukan Sekarang", "Jadwalkan", "Delegasikan", "Tinggalkan"]
    selected_priorities = st.multiselect("Filter berdasarkan prioritas:", options=priority_options)
    filtered_tasks = st.session_state.tasks
    if search_query: filtered_tasks = [task for task in filtered_tasks if search_query.lower() in task.get('Tugas', '').lower()]
    if selected_priorities: filtered_tasks = [task for task in filtered_tasks if task.get('Prioritas') in selected_priorities]
    if not filtered_tasks: st.warning("Tidak ada tugas yang cocok dengan kriteria pencarian Anda.")
    else:
        header_cols = st.columns((2.5, 3, 1.5, 1.5, 1.5, 1.5, 2))
        col_names = ["Tugas", "Deskripsi", "Status", "Tgl Mulai", "Tgl Selesai", "Prioritas", "Aksi"]
        for col, name in zip(header_cols, col_names): col.markdown(f'<p class="table-header">{name}</p>', unsafe_allow_html=True)
        for task in filtered_tasks:
            st.markdown("---")
            row_cols = st.columns((2.5, 3, 1.5, 1.5, 1.5, 1.5, 2))
            status_text, _ = get_task_status(task)
            row_cols[0].write(task.get('Tugas')); row_cols[1].write(task.get('Deskripsi'))
            row_cols[2].markdown(f'<div class="status-badge {get_task_status(task)[1]}">{status_text}</div>', unsafe_allow_html=True)
            row_cols[3].write(task.get('Tanggal Mulai').strftime('%d %b %Y') if task.get('Tanggal Mulai') else "N/A")
            row_cols[4].write(task.get('Tanggal Selesai Target').strftime('%d %b %Y') if task.get('Tanggal Selesai Target') else "N/A")
            row_cols[5].write(task.get('Prioritas'))
            with row_cols[6]:
                action_cols = st.columns(2)
                if action_cols[0].button("Edit", key=f"edit_{task['id']}"): st.session_state.editing_task_id = task['id']; st.rerun()
                if action_cols[1].button("Hapus", key=f"del_{task['id']}", type="primary"):
                    delete_task_from_db(task['id'])
                    st.session_state.tasks = [t for t in st.session_state.tasks if t['id'] != task['id']]
                    st.success(f"Tugas '{task.get('Tugas')}' dihapus."); st.rerun()

# --- KONTEN TAB 2: PRIORITASKAN ---
with tab2:
    st.header("Prioritaskan Tugas Anda")
    st.subheader("Daftar Tugas Aktif")
    if not tasks_aktif:
        st.info("Tidak ada tugas aktif untuk diprioritaskan.")
    else:
        df_aktif = pd.DataFrame(tasks_aktif)
        st.dataframe(
            df_aktif[['Tugas', 'Prioritas', 'Tanggal Mulai', 'Tanggal Selesai Target', 'Deskripsi']],
            use_container_width=True, hide_index=True
        )
    
    st.markdown("---")
    st.info("Pilih tugas dari 'Bank Tugas' di bawah dan klik tombol kuadran tujuan untuk memindahkannya.")
    
    bank_tugas = [task for task in tasks_aktif if task.get('Prioritas') in ['Belum Diprioritaskan', 'Bank Tugas']]
    if bank_tugas:
        selected_task_name = st.selectbox("Pilih Tugas dari Bank Tugas:", options=[t['Tugas'] for t in bank_tugas], label_visibility="collapsed")
        selected_task = next((t for t in bank_tugas if t.get('Tugas') == selected_task_name), None)
        if selected_task:
            st.subheader("Pindahkan Tugas Terpilih Ke:")
            prio_cols = st.columns(4)
            if prio_cols[0].button("Lakukan Sekarang"): update_task_priority_by_id(selected_task['id'], "Lakukan Sekarang"); st.rerun()
            if prio_cols[1].button("Jadwalkan"): update_task_priority_by_id(selected_task['id'], "Jadwalkan"); st.rerun()
            if prio_cols[2].button("Delegasikan"): update_task_priority_by_id(selected_task['id'], "Delegasikan"); st.rerun()
            if prio_cols[3].button("Tinggalkan"): update_task_priority_by_id(selected_task['id'], "Tinggalkan"); st.rerun()
    else:
        st.success("Semua tugas aktif sudah diprioritaskan!")
    
    st.markdown("---")
    
    tasks_needing_schedule = [task for task in tasks_aktif if task.get('Prioritas') == 'Jadwalkan' and not task.get('Tanggal Jadwal')]
    if tasks_needing_schedule:
        st.warning("Beberapa tugas perlu diatur tanggal eksekusinya:")
        with st.form("set_schedule_form"):
            for task in tasks_needing_schedule:
                # --- PERBAIKAN DI SINI ---
                # Gunakan tanggal tugas yang ada atau hari ini sebagai default
                default_schedule_date = task.get('Tanggal Jadwal') or date.today()
                st.date_input(
                    f"Tanggal eksekusi untuk '{task.get('Tugas')}'",
                    value=default_schedule_date, # Menggunakan nilai default yang valid
                    min_value=date.today(),
                    key=f"date_in_{task['id']}"
                )
                # --- AKHIR PERBAIKAN ---

            if st.form_submit_button("Simpan Semua Jadwal"):
                for task in tasks_needing_schedule:
                    task['Tanggal Jadwal'] = st.session_state[f"date_in_{task['id']}"]
                    save_task_to_db(task)
                st.success("Tanggal eksekusi berhasil disimpan!"); st.rerun()

with tab3:
    st.header("Jadwal Aktivitas Anda")
    tasks_to_schedule = sorted([t for t in tasks_aktif if t.get('Prioritas') in ["Lakukan Sekarang", "Jadwalkan"]], key=lambda x: x.get('Tanggal Jadwal') or x.get('Tanggal Mulai'))
    if not tasks_to_schedule: st.info("Tidak ada tugas yang perlu dijadwalkan.")
    else:
        st.success("Tandai tugas yang telah selesai di sini untuk mencatat kemajuan Anda.")
        for group_date, tasks_group in groupby(tasks_to_schedule, key=lambda x: x.get('Tanggal Jadwal') or x.get('Tanggal Mulai')):
            if group_date:
                nice_date = group_date.strftime("%A, %d %B %Y")
                header = f"Hari Ini: {nice_date}" if group_date == date.today() else f"{nice_date}"
                with st.container(border=True):
                    st.subheader(header)
                    for task in tasks_group:
                        c1, c2 = st.columns([0.1, 4])
                        with c1:
                            is_done = st.checkbox("", value=task.get('Selesai'), key=f"cb_final_{task['id']}", label_visibility="collapsed")
                            if is_done != task.get('Selesai'):
                                task['Selesai'] = is_done
                                task['Tanggal Selesai'] = date.today() if is_done else None
                                save_task_to_db(task); st.rerun()
                        with c2: st.write(f"{task.get('Tugas')} (Prioritas: {task.get('Prioritas')})")

with tab4:
    st.header("Visualisasi Kalender Tugas")
    calendar_events = []
    for task in st.session_state.tasks:
        start_date = task.get("Tanggal Mulai")
        end_date = task.get("Tanggal Selesai Target")
        if start_date and not end_date: end_date = start_date
        elif not start_date and end_date: start_date = end_date
        if start_date and end_date:
            end_date_for_calendar = end_date + timedelta(days=1)
            _, color = get_task_status(task)
            calendar_events.append({"title": task["Tugas"], "start": start_date.isoformat(), "end": end_date_for_calendar.isoformat(), "color": color})
    calendar_options = {"headerToolbar": {"left": "prev,next today", "center": "title", "right": "dayGridMonth,timeGridWeek,timeGridDay"}, "initialView": "dayGridMonth", "height": "700px"}
    if calendar_events: calendar(events=calendar_events, options=calendar_options, key="calendar")
    else: st.info("Tidak ada tugas dengan tanggal mulai/selesai yang bisa ditampilkan di kalender.")

with tab5:
    st.header("Sesi Fokus (Pomodoro)")
    POMODORO_DURATION = 25 * 60
    if st.session_state.pomodoro_running:
        elapsed_time = time.time() - st.session_state.pomodoro_start_time
        time_left = POMODORO_DURATION - elapsed_time
        if time_left > 0:
            st.warning(f"Sesi fokus untuk: **{st.session_state.active_pomodoro_task}**")
            if st.button("Hentikan Sesi"): st.session_state.pomodoro_running = False; st.session_state.active_pomodoro_task = None; st.rerun()
            st.progress(elapsed_time / POMODORO_DURATION)
            minutes, seconds = divmod(int(time_left), 60)
            st.metric("Sisa Waktu Kerja:", f"{minutes:02d}:{seconds:02d}")
            time.sleep(1); st.rerun()
        else:
            st.success(f"Sesi fokus untuk '{st.session_state.active_pomodoro_task}' telah selesai!")
            st.balloons(); st.session_state.pomodoro_running = False; st.session_state.active_pomodoro_task = None
            time.sleep(3); st.rerun()
    else:
        st.info("Pilih tugas dengan prioritas 'Lakukan Sekarang' di bawah ini untuk memulai sesi Pomodoro.")
        tasks_to_focus = [t for t in tasks_aktif if t.get('Prioritas') == "Lakukan Sekarang"]
        if not tasks_to_focus: st.warning("Tidak ada tugas dengan prioritas 'Lakukan Sekarang' yang tersedia.")
        else:
            for task in tasks_to_focus:
                 c1, c2 = st.columns([3, 1])
                 with c1: st.write(task.get('Tugas'))
                 with c2:
                    if st.button("Mulai Fokus", key=f"focus_center_{task['id']}"):
                        st.session_state.pomodoro_running = True; st.session_state.active_pomodoro_task = task.get('Tugas')
                        st.session_state.pomodoro_start_time = time.time(); st.rerun()

# --- (DIUBAH) TAB 7: LAPORAN & EKSPOR ---
with tabs[6]:
    st.header("Laporan & Ekspor")

    st.subheader("Analisis Produktivitas")
    if not st.session_state.tasks:
        st.info("Belum ada data tugas untuk ditampilkan dalam laporan.")
    else:
        df_all_tasks = pd.DataFrame(st.session_state.tasks)
        st.markdown("#### Tugas Selesai dalam 7 Hari Terakhir")
        if 'Tanggal Selesai' in df_all_tasks.columns and not df_all_tasks['Tanggal Selesai'].isnull().all():
            df_all_tasks['Tanggal Selesai'] = pd.to_datetime(df_all_tasks['Tanggal Selesai'], errors='coerce').dt.date
            df_all_tasks.dropna(subset=['Tanggal Selesai'], inplace=True)
            today = date.today()
            seven_days_ago = today - timedelta(days=6)
            completed_tasks_recent = df_all_tasks[(df_all_tasks['Selesai'] == True) & (df_all_tasks['Tanggal Selesai'] >= seven_days_ago) & (df_all_tasks['Tanggal Selesai'] <= today)]
            if completed_tasks_recent.empty:
                st.warning("Tidak ada tugas yang diselesaikan dalam 7 hari terakhir.")
            else:
                tasks_per_day = completed_tasks_recent.groupby('Tanggal Selesai').size().reset_index(name='Jumlah')
                date_range = pd.date_range(start=seven_days_ago, end=today).to_frame(index=False, name='Tanggal Selesai')
                date_range['Tanggal Selesai'] = date_range['Tanggal Selesai'].dt.date
                merged_data = pd.merge(date_range, tasks_per_day, on='Tanggal Selesai', how='left').fillna(0)
                fig_bar = px.bar(merged_data, x='Tanggal Selesai', y='Jumlah', title='Jumlah Tugas Selesai per Hari', labels={'Tanggal Selesai': 'Tanggal', 'Jumlah': 'Jumlah Tugas'}, text_auto=True)
                fig_bar.update_traces(marker_color='#007BFF')
                st.plotly_chart(fig_bar, use_container_width=True)
        else:
            st.warning("Tidak ada data tugas selesai untuk ditampilkan.")
        st.markdown("---")
        st.markdown("#### Distribusi Prioritas Tugas Aktif")
        active_tasks_df = pd.DataFrame(tasks_aktif)
        if active_tasks_df.empty:
            st.success("Selamat! Tidak ada tugas aktif saat ini.")
        else:
            priority_counts = active_tasks_df['Prioritas'].value_counts().reset_index(name='Jumlah')
            fig_pie = px.pie(priority_counts, names='Prioritas', values='Jumlah', title='Komposisi Prioritas Tugas yang Sedang Berjalan', hole=0.3)
            fig_pie.update_traces(textposition='inside', textinfo='percent+label')
            st.plotly_chart(fig_pie, use_container_width=True)

    st.markdown("---")
    # --- (BARU) BAGIAN EKSPOR PPT ---
    st.subheader("Ekspor ke PowerPoint")
    st.info("Buat laporan ringkas dalam format presentasi PowerPoint (.pptx) berdasarkan data tugas Anda saat ini.")
    
    pptx_buffer = buat_presentasi_laporan(st.session_state.tasks)
    
    st.download_button(
        label="Download Laporan (.pptx)",
        data=pptx_buffer,
        file_name=f"Laporan_Produktivitas_{date.today()}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )